import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- COLORES CORPORATIVOS ---
COLOR_PRIMARY = RGBColor(0, 51, 102)    # Azul Oscuro
COLOR_ACCENT = RGBColor(255, 102, 0)    # Naranja
COLOR_GRAY = RGBColor(100, 100, 100)    # Gris Texto
COLOR_LIGHT = RGBColor(245, 245, 245)   # Gris Fondo

def create_pptx_from_structure(data):
    """
    Generador Inteligente de Slides. Detecta la estructura del JSON y 
    elige el diseño visual adecuado (Cuadrantes, Tablas, Columnas).
    """
    prs = Presentation()
    
    # 1. Forzar formato Widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 2. Crear Slide en blanco
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)
    
    # 3. Dibujar Cabecera (Título y Subtítulo)
    _draw_header(slide, data)
    
    # 4. ENRUTADOR INTELIGENTE
    # Detectamos qué tipo de información trajo la IA
    keys = " ".join(data.keys()).lower()
    t_type = data.get("template_type", "").lower()
    
    if "dofa" in t_type or "swot" in t_type or ("fortalezas" in keys and "amenazas" in keys):
        _draw_dofa_layout(slide, data)
        
    elif "buyer" in t_type or "persona" in t_type or ("perfil" in keys and "dolor" in keys):
        _draw_persona_layout(slide, data)
        
    elif "journey" in t_type or "viaje" in t_type or "etapa" in keys:
        _draw_journey_table(slide, data)
        
    elif "matriz" in t_type or "2x2" in t_type:
        _draw_matrix_layout(slide, data)
        
    else:
        # Si no es ninguno específico, usamos el genérico inteligente
        # que imprime TODO lo que encuentre.
        _draw_smart_generic_layout(slide, data)

    # 5. Insight al pie de página (si existe)
    if data.get("insight_principal"):
        _draw_footer_insight(slide, data.get("insight_principal"))

    # Retornar archivo
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ==========================================
# FUNCIONES DE DIBUJO ESPECÍFICAS
# ==========================================

def _draw_header(slide, data):
    # Título
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = str(data.get("titulo", "Sin Título")).upper()
    p.font.name = 'Arial'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Subtítulo
    if data.get("subtitulo"):
        tb_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.5))
        p_sub = tb_sub.text_frame.paragraphs[0]
        p_sub.text = str(data.get("subtitulo", ""))
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = COLOR_GRAY

def _draw_dofa_layout(slide, data):
    """Dibuja una matriz 2x2 para Fortalezas, Oportunidades, Debilidades, Amenazas."""
    # Coordenadas base
    margin_x = 0.5
    margin_y = 1.8
    w = 6.0
    h = 2.2
    gap = 0.2
    
    # Mapeo de cuadrantes
    quadrants = [
        ("Fortalezas", data.get("fortalezas", []), margin_x, margin_y),
        ("Debilidades", data.get("debilidades", []), margin_x + w + gap, margin_y),
        ("Oportunidades", data.get("oportunidades", []), margin_x, margin_y + h + gap),
        ("Amenazas", data.get("amenazas", []), margin_x + w + gap, margin_y + h + gap)
    ]
    
    for title, points, x, y in quadrants:
        # Fondo
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_LIGHT
        shape.line.color.rgb = COLOR_PRIMARY
        
        # Título del cuadrante
        tb = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.1), Inches(w-0.2), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = title.upper()
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p.font.size = Pt(14)
        
        # Contenido (Bullets)
        tb_body = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.5), Inches(w-0.2), Inches(h-0.6))
        tf = tb_body.text_frame
        tf.word_wrap = True
        
        points_list = points if isinstance(points, list) else [str(points)]
        for point in points_list[:5]: # Limitar a 5 bullets para que quepa
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(11)

def _draw_journey_table(slide, data):
    """Detecta etapas y crea una tabla visual."""
    # Buscar claves que parezcan etapas (etapa_1, etapa_2...)
    stages = [v for k, v in data.items() if "etapa" in k.lower() and isinstance(v, dict)]
    
    if not stages:
        _draw_smart_generic_layout(slide, data)
        return

    cols = len(stages)
    rows = 3 # Nombre, Acción, Pensamiento
    
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.0))
    table = table_shape.table
    
    for i, stage in enumerate(stages):
        # Fila 1: Header (Nombre Etapa)
        cell = table.cell(0, i)
        cell.text = str(stage.get("nombre", f"Etapa {i+1}")).upper()
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Fila 2: Acciones
        cell = table.cell(1, i)
        act = stage.get("accion", stage.get("actividades", ""))
        cell.text = f"Acción:\n{act}"
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Fila 3: Pensamientos/Sentimientos
        cell = table.cell(2, i)
        th = stage.get("pensamiento", stage.get("sentimiento", ""))
        cell.text = f"Piensa/Siente:\n{th}"
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].font.italic = True

def _draw_persona_layout(slide, data):
    """Diseño con barra lateral para perfil."""
    # 1. Sidebar (Datos demográficos)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(3.5), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT
    
    tf = shape.text_frame
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = str(data.get("perfil_nombre", "Usuario")).upper()
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.font.size = Pt(16)
    
    details = [
        f"Bio: {data.get('perfil_demografia', '')}",
        f"Edad: {data.get('edad', '')}",
        f"Ocupación: {data.get('ocupacion', '')}"
    ]
    for d in details:
        p = tf.add_paragraph()
        p.text = d
        p.font.size = Pt(11)
        p.space_after = Pt(6)
        
    # 2. Área principal (Necesidades y Dolores)
    # Dividimos el resto en 2 cajas verticales
    # Caja Superior: Necesidades
    tb1 = slide.shapes.add_textbox(Inches(4.2), Inches(1.8), Inches(8.5), Inches(2.0))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = "NECESIDADES & MOTIVACIONES"
    p1.font.bold = True
    p1.font.color.rgb = COLOR_ACCENT
    
    needs = data.get("necesidades_jtbd", []) + data.get("deseos_motivaciones", [])
    for n in needs[:4]:
        p = tb1.text_frame.add_paragraph()
        p.text = f"• {n}"
        
    # Caja Inferior: Dolores
    tb2 = slide.shapes.add_textbox(Inches(4.2), Inches(4.0), Inches(8.5), Inches(2.0))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "FRUSTRACIONES & DOLORES"
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ACCENT
    
    pains = data.get("puntos_dolor_frustraciones", [])
    for pn in pains[:4]:
        p = tb2.text_frame.add_paragraph()
        p.text = f"• {pn}"

def _draw_smart_generic_layout(slide, data):
    """
    Intenta dibujar TODO lo que no sea título/subtítulo en una lista limpia.
    Evita que salgan slides vacíos.
    """
    ignore_keys = ["titulo", "subtitulo", "insight_principal", "template_type", "titulo_diapositiva"]
    
    content_keys = [k for k in data.keys() if k not in ignore_keys]
    
    # Caja de texto grande
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    first = True
    for key in content_keys:
        val = data[key]
        if not val: continue
        
        # Título de la sección (ej: "Puntos Clave")
        if not first: p = tf.add_paragraph()
        else: p = tf.paragraphs[0]
        
        clean_key = key.replace("_", " ").upper()
        p.text = clean_key
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.font.size = Pt(12)
        p.space_before = Pt(12)
        
        # Contenido
        items = val if isinstance(val, list) else [str(val)]
        for item in items:
            bullet = tf.add_paragraph()
            bullet.text = f"• {item}"
            bullet.font.size = Pt(11)
            bullet.space_after = Pt(2)
            
        first = False

def _draw_footer_insight(slide, text):
    # Caja destacada al final
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 240, 230) # Naranja muy suave
    shape.line.color.rgb = COLOR_ACCENT
    
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"💡 INSIGHT: {text}"
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER
    
    # Hack para centrar verticalmente en python-pptx básico
    # (El ajuste fino de márgenes ayuda)
    tf.margin_top = Inches(0.2)
