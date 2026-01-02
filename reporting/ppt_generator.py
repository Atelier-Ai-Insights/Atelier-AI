from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import io
import re
import pandas as pd

def crear_ppt_desde_json(data_json):
    """
    Genera un PowerPoint ajustado EXACTAMENTE a 40.64cm x 22.86cm (16x9 pulgadas).
    Recalibrado para llenar todo el lienzo.
    """
    
    # 1. Cargar Plantilla Base
    try:
        prs = Presentation("Plantilla_PPT_ATL.pptx")
    except:
        prs = Presentation()

    # --- CONFIGURACIÓN EXACTA SEGÚN TU IMAGEN ---
    prs.slide_width = Inches(16)  # 40.64 cm
    prs.slide_height = Inches(9)  # 22.86 cm

    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(slide_layout)

    # 2. Configurar Título (Centrado en 16 pulgadas)
    titulo_texto = data_json.get('titulo_diapositiva', 'Resumen Estratégico')
    
    if slide.shapes.title:
        slide.shapes.title.text = titulo_texto
    else:
        # Caja de título ancha (15 pulgadas)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15.0), Inches(1.2))
        tf = title_box.text_frame
        tf.text = titulo_texto
        tf.paragraphs[0].font.size = Pt(32) # Fuente más grande para lienzo grande
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 3. Router Inteligente
    template_type = data_json.get('template_type', '').lower()
    keys_str = " ".join(data_json.keys()).lower() 
    
    if "matriz" in template_type or "2x2" in template_type or "cuadrante" in keys_str:
        _dibujar_matriz_nativa(slide, data_json)

    elif "foda" in template_type or "swot" in template_type or "dofa" in template_type or ("fortalezas" in keys_str and "amenazas" in keys_str):
        _dibujar_foda_nativo(slide, data_json)

    elif "embudo" in template_type or "funnel" in template_type or "conversion" in keys_str:
        _dibujar_embudo_nativo(slide, data_json)

    elif "journey" in template_type or "viaje" in template_type or "map" in template_type or "etapa 1" in keys_str:
        _dibujar_journey_nativo(slide, data_json)

    elif "persona" in template_type or "buyer" in template_type or "perfil" in template_type or ("demografia" in keys_str and "frustraciones" in keys_str):
        _dibujar_buyer_persona_nativo(slide, data_json)

    elif "empatia" in template_type or "empathy" in template_type or ("dice" in keys_str and "piensa" in keys_str):
        _dibujar_mapa_empatia_nativo(slide, data_json)

    elif "valor" in template_type or "value" in template_type or ("alegrias" in keys_str and "dolores" in keys_str):
        _dibujar_propuesta_valor_nativo(slide, data_json)

    else:
        _dibujar_lista_generica(slide, data_json)

    # 4. Conclusión (Ancho completo 15 pulgadas)
    if 'conclusion_clave' in data_json:
        # Pie de página ancho
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(7.8), Inches(15.0), Inches(1.0))
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(245, 245, 245); bg.line.color.rgb = RGBColor(220, 220, 220)
        tf = bg.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
        tf.text = "💡 " + data_json['conclusion_clave']
        p = tf.paragraphs[0]; p.font.color.rgb = RGBColor(50, 50, 50); p.alignment = PP_ALIGN.LEFT; p.font.size = Pt(14)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

def add_analysis_slide(prs, type, title, content):
    """
    Helper unificado para añadir slides de Análisis de Datos (Tablas, Imágenes).
    Integrado en Fase 3 para limpiar data_analysis_mode.py.
    """
    try:
        # Layout 5 suele ser Titulo + Contenido en blanco o Titulo + Objeto
        # Ajustar según tu plantilla. Si falla, usa el último disponible.
        layout_index = 5 if len(prs.slide_layouts) > 5 else -1
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        if type == "image":
            if content is None: return
            # Resetear puntero del buffer de imagen
            content.seek(0)
            # Centrar imagen aproximadamente
            slide.shapes.add_picture(content, Inches(1.0), Inches(2.0), height=Inches(5.5))
        
        elif type == "table":
            if content is None or content.empty: return
            
            # Aplanar índice si es MultiIndex
            df = content.reset_index() if (content.index.name or isinstance(content.index, pd.MultiIndex)) else content
            
            rows, cols = df.shape
            # Limite de seguridad
            if rows > 12: df = df.head(12); rows = 12
            
            # Crear tabla
            graphic_frame = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(2.0), Inches(15.0), Inches(5.0))
            table = graphic_frame.table
            
            # Headers
            for c in range(cols):
                cell = table.cell(0, c)
                cell.text = str(df.columns[c])
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
                if cell.text_frame.paragraphs:
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.bold = True

            # Body
            for r in range(rows):
                for c in range(cols):
                    val = df.iloc[r, c]
                    cell = table.cell(r+1, c)
                    cell.text = f"{val:.2f}" if isinstance(val, (float, int)) else str(val)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    
    except Exception as e:
        print(f"Error generando slide tipo {type}: {e}")


# ==============================================================================
# FUNCIONES DE DIBUJO (MANTENIDAS DE LA VERSIÓN ANTERIOR)
# ==============================================================================
# ... (Aquí va el resto de funciones privadas _dibujar_buyer_persona_nativo, etc. 
# ...  MANTÉN EL CÓDIGO ORIGINAL DE ESAS FUNCIONES, NO CAMBIAN).
# ... Solo asegúrate de copiar las funciones _dibujar_* y los helpers _buscar_clave_flexible, etc.
# ... Para brevedad, asumo que mantienes esas funciones auxiliares abajo.

def _dibujar_buyer_persona_nativo(slide, data):
    """Diseño 16:9 - Sidebar Izquierda + Paneles Anchos."""
    # (Mantener código original...)
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(3.5), Inches(6.0))
    sidebar.fill.solid(); sidebar.fill.fore_color.rgb = RGBColor(230, 240, 250); sidebar.line.fill.background()
    # ... (Resto de la función original)

def _dibujar_matriz_nativa(slide, data):
    # (Mantener código original...)
    pass 

def _dibujar_foda_nativo(slide, data):
    # (Mantener código original...)
    pass

def _dibujar_journey_nativo(slide, data):
    # (Mantener código original...)
    pass

def _dibujar_mapa_empatia_nativo(slide, data):
    # (Mantener código original...)
    pass

def _dibujar_propuesta_valor_nativo(slide, data):
    # (Mantener código original...)
    pass

def _dibujar_embudo_nativo(slide, data):
    # (Mantener código original...)
    pass

def _dibujar_lista_generica(slide, data):
    # (Mantener código original...)
    left = Inches(0.5); top = Inches(1.5); width = Inches(15.0); height = Inches(6.0)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    excluded_keys = ['titulo_diapositiva', 'template_type', 'conclusion_clave']
    first = True
    for k, v in data.items():
        if k in excluded_keys: continue
        if not first: tf.add_paragraph()
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        p.text = k.replace('_', ' ').upper(); p.font.bold = True; p.font.color.rgb = RGBColor(0, 51, 102); p.font.size = Pt(14); first = False
        _llenar_text_frame_flexible(tf, v if isinstance(v, list) else [v])

# Helpers (Importante mantenerlos)
def _buscar_clave_flexible(data, lista_keywords):
    for kw in lista_keywords:
        if kw in data: return data[kw]
    for key_json, val in data.items():
        key_clean = key_json.lower()
        for kw in lista_keywords:
            if kw.lower() in key_clean: return val
    return None

def _llenar_text_frame_flexible(text_frame, lista_items):
    if not lista_items: return
    if not isinstance(lista_items, list): lista_items = [str(lista_items)]
    valid_items = [str(x) for x in lista_items if x and str(x).lower() != "none"]
    if not valid_items: return
    p = text_frame.paragraphs[0]; p.text = f"• {valid_items[0]}"; p.font.color.rgb = RGBColor(40, 40, 40)
    for item in valid_items[1:]:
        p = text_frame.add_paragraph(); p.text = f"• {item}"; p.font.color.rgb = RGBColor(40, 40, 40)

def _llenar_text_frame_tabla(text_frame, content):
    text_frame.clear()
    if content is None or str(content).strip().lower() == "none": content = "-"
    if isinstance(content, list):
        for item in content:
            if not item: continue
            p = text_frame.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(10); p.font.color.rgb = RGBColor(0, 0, 0); p.space_after = Pt(2)
    else:
        text_str = str(content).strip()
        if not text_str: return
        p = text_frame.add_paragraph(); p.text = text_str; p.font.size = Pt(10); p.font.color.rgb = RGBColor(0, 0, 0)

def _crear_etiqueta(slide, x, y, texto, bold=False, vertical=False):
    w, h = (Inches(2), Inches(0.5)) if not vertical else (Inches(0.5), Inches(2))
    tb = slide.shapes.add_textbox(Inches(x) - w/2, Inches(y) - h/2, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.text = str(texto); p.alignment = PP_ALIGN.CENTER; p.font.bold = bold; p.font.color.rgb = RGBColor(80, 80, 80); p.font.size = Pt(12)
    if vertical: tb.rotation = -90

def _poner_titulo_contenido(slide, shape, title, content):
    tf = shape.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.text = title.upper(); p.font.bold = True; p.font.size = Pt(11); p.alignment = PP_ALIGN.CENTER
    tf.add_paragraph()
    _llenar_text_frame_flexible(tf, content)

def _poner_titulo_contenido_manual(slide, x, y, w, h, title, content):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.text = title.upper(); p.font.bold = True; p.font.size = Pt(11); p.alignment = PP_ALIGN.CENTER
    tf.add_paragraph()
    _llenar_text_frame_flexible(tf, content)
